# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from textwrap import fill

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT / "output"
OUTPUT_PPTX = OUTPUT_DIR / "HRD_AI_Workshop_Structured_Deck_20260417.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG_LIGHT = "F5F1EA"
BG_DARK = "10283C"
NAVY = "16384C"
TEAL = "236D73"
SAGE = "6A8A7A"
ACCENT = "C86E4E"
GOLD = "C9A25D"
CARD = "FFFDFC"
BORDER = "D9D0C3"
TEXT = "1F2E38"
MUTED = "61717B"
WHITE = "FFFFFF"

TITLE_FONT = "Malgun Gothic"
BODY_FONT = "Malgun Gothic"
MONO_FONT = "Consolas"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def i(value: float):
    return Inches(value)


def pt(value: float):
    return Pt(value)


def wrap_file(text: str, width: int = 26) -> str:
    tokens = (
        text.replace("/", "/\n")
        .replace(" (`", "\n(`")
        .replace("_", "_\n")
        .splitlines()
    )
    wrapped = []
    for token in tokens:
        if len(token) <= width:
            wrapped.append(token)
        else:
            wrapped.append(fill(token, width=width, break_long_words=True))
    return "\n".join(part for part in wrapped if part)


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_shape(slide, shape_type, x, y, w, h, fill_color, line_color=None, line_width=1.0):
    shape = slide.shapes.add_shape(shape_type, i(x), i(y), i(w), i(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    *,
    font_name=BODY_FONT,
    font_size=18,
    color=TEXT,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_left=6,
    margin_right=6,
    margin_top=4,
    margin_bottom=4,
):
    box = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = pt(margin_left)
    tf.margin_right = pt(margin_right)
    tf.margin_top = pt(margin_top)
    tf.margin_bottom = pt(margin_bottom)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = pt(font_size)
    font.color.rgb = rgb(color)
    font.bold = bold
    font.italic = italic
    return box


def add_paragraphs(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    *,
    font_name=BODY_FONT,
    font_size=15,
    color=TEXT,
    bullet_prefix="• ",
    level_indent=0.18,
    line_spacing=1.12,
    space_after=4,
    margin_left=6,
    margin_right=6,
    margin_top=4,
    margin_bottom=4,
):
    box = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = pt(margin_left)
    tf.margin_right = pt(margin_right)
    tf.margin_top = pt(margin_top)
    tf.margin_bottom = pt(margin_bottom)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    normalized = []
    for item in paragraphs:
        if isinstance(item, tuple):
            normalized.append({"text": item[1], "level": item[0]})
        elif isinstance(item, dict):
            normalized.append(item)
        else:
            normalized.append({"text": str(item), "level": 0})

    for idx, item in enumerate(normalized):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        level = item.get("level", 0)
        prefix = ("   " * level) + (item.get("prefix") or ("◦ " if level else bullet_prefix))
        run = p.add_run()
        run.text = f"{prefix}{item['text']}"
        font = run.font
        font.name = item.get("font_name", font_name)
        font.size = pt(item.get("font_size", font_size))
        font.bold = item.get("bold", False)
        font.italic = item.get("italic", False)
        font.color.rgb = rgb(item.get("color", color))
        p.line_spacing = line_spacing
        p.space_after = pt(item.get("space_after", space_after))
        p.left_indent = i(level_indent * level)
    return box


def add_slide_number(slide, number: int, *, color=MUTED) -> None:
    add_textbox(
        slide,
        12.55,
        7.0,
        0.45,
        0.25,
        str(number),
        font_name=BODY_FONT,
        font_size=10,
        color=color,
        align=PP_ALIGN.RIGHT,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_section_label(slide, label: str, *, color=TEAL, text_color=WHITE) -> None:
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.65, 0.35, 1.42, 0.32, color, None)
    add_textbox(
        slide,
        0.72,
        0.39,
        1.28,
        0.22,
        label,
        font_size=10,
        color=text_color,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_content_title(slide, label: str, title: str, subtitle: str | None = None, number: int | None = None) -> None:
    add_section_label(slide, label)
    add_textbox(
        slide,
        0.65,
        0.78,
        9.8,
        0.48,
        title,
        font_name=TITLE_FONT,
        font_size=28,
        color=TEXT,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    if subtitle:
        add_textbox(
            slide,
            0.68,
            1.22,
            10.0,
            0.28,
            subtitle,
            font_name=BODY_FONT,
            font_size=11,
            color=MUTED,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.84, 0.22, 1.28, 0.38, ACCENT, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.18, 0.58, 1.95, 0.16, GOLD, None)
    if number is not None:
        add_slide_number(slide, number)


def add_card(slide, x, y, w, h, title: str, body=None, *, accent=TEAL, body_font=14, title_font=16):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, CARD, BORDER, 1.0)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.14, h, accent, None)
    add_textbox(
        slide,
        x + 0.24,
        y + 0.16,
        w - 0.38,
        0.35,
        title,
        font_name=TITLE_FONT,
        font_size=title_font,
        color=TEXT,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    if isinstance(body, str) and body:
        add_textbox(
            slide,
            x + 0.22,
            y + 0.58,
            w - 0.34,
            h - 0.68,
            body,
            font_name=BODY_FONT,
            font_size=body_font,
            color=TEXT,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )
    elif body:
        add_paragraphs(
            slide,
            x + 0.18,
            y + 0.52,
            w - 0.3,
            h - 0.62,
            body,
            font_size=body_font,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )


def add_divider_slide(prs: Presentation, section_no: str, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_DARK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.0, 2.2, 7.5, TEAL, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.1, 0.0, 1.23, 7.5, ACCENT, None)
    add_textbox(
        slide,
        0.78,
        1.2,
        2.8,
        0.46,
        f"SECTION {section_no}",
        font_name=BODY_FONT,
        font_size=14,
        color=WHITE,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(
        slide,
        0.82,
        2.0,
        3.6,
        1.2,
        section_no,
        font_name=TITLE_FONT,
        font_size=94,
        color=GOLD,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(
        slide,
        3.9,
        2.28,
        7.5,
        0.72,
        title,
        font_name=TITLE_FONT,
        font_size=30,
        color=WHITE,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(
        slide,
        3.92,
        3.12,
        6.9,
        1.18,
        subtitle,
        font_name=BODY_FONT,
        font_size=18,
        color="DDE6EA",
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.9, 4.62, 4.15, 0.6, CARD, None)
    add_textbox(
        slide,
        4.12,
        4.82,
        3.6,
        0.2,
        "구조화된 흐름으로 다음 슬라이드를 전개",
        font_name=BODY_FONT,
        font_size=12,
        color=TEXT,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_DARK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 0.0, 13.33, 0.26, GOLD, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.4, 0.26, 2.93, 7.24, ACCENT, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.85, 0.78, 3.3, 5.58, CARD, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.75, 0.9, 0.16, 3.72, TEAL, None)
    add_textbox(
        slide,
        1.08,
        1.15,
        7.2,
        1.0,
        "[실습] HRD부문의\n효과적인 AI 활용 방안",
        font_name=TITLE_FONT,
        font_size=30,
        color=WHITE,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(
        slide,
        1.1,
        2.35,
        6.6,
        0.4,
        "부제: (주)현대제철 사례를 중심으로",
        font_name=BODY_FONT,
        font_size=17,
        color="DDE6EA",
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(
        slide,
        1.12,
        3.0,
        6.4,
        1.25,
        "강의 메모 수준의 초안을 실제 교육 운영안과 발표 설계 문서 수준의 발표자료로 구조화",
        font_name=BODY_FONT,
        font_size=19,
        color=WHITE,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )
    add_textbox(slide, 9.08, 1.12, 2.5, 0.25, "Workshop Meta", font_size=12, color=TEAL, bold=True, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    meta = [
        "일시\n2026.04.17 (금)\n10:00~17:00",
        "주관\n(주)현대제철\n컬처디자인팀",
        "강사\n한충석\n책임매니저",
        "참석 명단\n`DOCUMENT_4월 HRD부문의 AI 활용 세미나 (1차명단).pdf` 참조",
    ]
    y = 1.55
    for item in meta:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.02, y, 2.6, 1.0, BG_LIGHT, BORDER, 0.8)
        add_textbox(
            slide,
            9.18,
            y + 0.12,
            2.24,
            0.74,
            item,
            font_name=BODY_FONT,
            font_size=12,
            color=TEXT,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )
        y += 1.12
    add_textbox(
        slide,
        1.1,
        6.85,
        5.3,
        0.22,
        "Source document: DOCUMENT_Design.md",
        font_name=BODY_FONT,
        font_size=10,
        color="C9D4DA",
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_overview_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "OVERVIEW", "과정 개요와 기대 산출물", "교육 대상, 목적, 운영 메모, 산출물을 한 장에서 조망", len(prs.slides))
    add_card(
        slide,
        0.72,
        1.65,
        3.0,
        2.08,
        "교육 대상",
        ["HRD 실무자", "교육 기획/운영 담당자", "조직문화/리더십 프로그램 담당자"],
        accent=TEAL,
    )
    add_card(
        slide,
        3.95,
        1.65,
        3.2,
        2.08,
        "교육 목적",
        [
            "생성형 AI의 기술 변화 방향 이해",
            "HRD 전 과정의 적용 시나리오 설계",
            "실무에 바로 연결되는 업무 문제 정의 역량 강화",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        7.38,
        1.65,
        2.9,
        2.08,
        "운영 메모",
        [
            "문서 목적: 메모 초안을 실제 운영안·발표 설계 문서로 정리",
            "참석 명단은 별도 PDF 참고",
            "샘플 마이크로사이트: 01chungee10snu.github.io/CL_Workshop",
        ],
        accent=SAGE,
        body_font=12,
    )
    add_card(
        slide,
        10.52,
        1.65,
        2.12,
        2.08,
        "핵심 산출물",
        [
            "GPTs 챗봇",
            "교육운영용 마이크로사이트 1종",
            "리더십 교육 F/U 개인화 메시지 설계",
            "만족도 데이터 자동화 분석 흐름도",
            "전사 AX 추진용 HRD 기획방안",
        ],
        accent=GOLD,
        body_font=11.5,
    )
    add_card(
        slide,
        0.72,
        4.0,
        11.92,
        2.46,
        "산출물 관점에서 본 이번 워크숍",
        [
            "기획: HRD 담당자가 업무 문제를 명확히 정의하고, 검증 가능한 결과물로 연결하는 연습",
            "운영: FAQ, 웹페이지, 후속안내, 실습자료처럼 반복되는 운영 업무를 구조화",
            "분석: 만족도 데이터의 정량·정성 분석을 자동화해 보고서 작성 시간을 단축",
            "확산: 2025 성과와 2026 운영방향, Change Agent 체계까지 하나의 흐름으로 연결",
        ],
        accent=NAVY,
        body_font=15,
        title_font=18,
    )


def add_core_message_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "MESSAGE", "핵심 메시지", "AI를 도구가 아니라 문제 정의와 산출물 연결의 가속기로 해석", len(prs.slides))
    messages = [
        ("01 생산성 증폭", "AI는 HRD 업무를 통째로 대체하기보다, 기획-운영-분석-커뮤니케이션의 생산성을 크게 증폭시킨다.", TEAL),
        ("02 HRD의 본질 유지", "HRD의 본질은 사람과 조직을 이해하고 행동 변화를 설계하는 일이며, AI는 그 설계를 더 빠르게 실험하고 구현하게 해주는 도구다.", ACCENT),
        ("03 문제 정의 역량", "중요한 것은 도구를 잘 사용하는 것이 아니라, 업무 문제를 명확히 정의하고 검증 가능한 산출물로 연결하는 능력이다.", NAVY),
    ]
    x = 0.82
    for title, body, color in messages:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.8, 4.0, 3.65, CARD, BORDER, 1.0)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.8, 4.0, 0.44, color, None)
        add_textbox(slide, x + 0.18, 2.0, 3.55, 0.28, title, font_size=15, color=WHITE, bold=True, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
        add_textbox(slide, x + 0.24, 2.52, 3.45, 2.32, body, font_size=20, color=TEXT, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
        x += 4.17
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.86, 5.86, 7.6, 0.72, BG_DARK, None)
    add_textbox(
        slide,
        3.08,
        6.08,
        7.18,
        0.24,
        "결론: '도구 사용법'보다 '업무 문제 정의 → 검증 가능한 산출물'의 흐름을 설계하는 것이 핵심",
        font_size=13,
        color=WHITE,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "AGENDA", "시간 배분(안)", "하루 전체 흐름을 세션 목적 중심으로 재구성", len(prs.slides))
    rows = [
        ("10:00~10:50", "강의", "생성형 AI 동향 이해, HRD 적용 관점 정리", NAVY),
        ("11:00~12:00", "실습 1", "교육운영 효율화를 위한 FAQ 챗봇 및 웹페이지 제작", TEAL),
        ("12:00~13:00", "점심", "휴식 및 질의응답", GOLD),
        ("13:00~13:50", "시연 1", "LLM API 활용 개인화 피드백 메시지 제작 및 개인화 메일 발송", ACCENT),
        ("14:00~15:30", "실습 2", "교육 데이터 분석 자동화 및 감정분석, 시각화", TEAL),
        ("15:40~16:30", "사례공유", "2025 성과와 2026 방향, Change Agent 운영", SAGE),
        ("16:30~17:00", "Wrap-up", "HRD 담당자의 AI 학습법과 전사 AX 제언", NAVY),
    ]
    y = 1.52
    for time, session, purpose, color in rows:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, y, 11.84, 0.7, CARD, BORDER, 1.0)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.02, y + 0.12, 1.7, 0.45, BG_DARK, None)
        add_textbox(slide, 1.18, y + 0.23, 1.35, 0.18, time, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.0, y + 0.12, 1.2, 0.45, color, None)
        add_textbox(slide, 3.12, y + 0.23, 0.95, 0.18, session, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
        add_textbox(slide, 4.46, y + 0.17, 7.92, 0.24, purpose, font_size=15, color=TEXT, bold=True, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
        y += 0.78


def add_lecture_goal_timeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "LECTURE", "강의 목표와 생성형 AI 발전 흐름", "학습목표와 기술 변화를 한 번에 파악하는 오프닝 슬라이드", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        3.5,
        4.95,
        "학습목표",
        [
            "생성형 AI의 발전 흐름 이해",
            "LLM, Transformer, Attention의 핵심 개념을 HRD 관점으로 단순화해 설명",
            "HRD 업무 전 과정에서 AI 적용 가능 지점을 구체적으로 식별",
        ],
        accent=TEAL,
        body_font=15,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 4.78, 2.0, 0.08, 3.95, BORDER, None)
    timeline = [
        ("2017", "Transformer 발표", "구글이 병렬 문맥 파악이 가능한 Transformer 아키텍처를 공개"),
        ("2020", "LLM 규모 경쟁", "GPT-3 등 거대언어모델의 파라미터 경쟁 본격화"),
        ("2022", "ChatGPT 대중화", "대화형 생성형 AI의 실무 도입과 대중화가 급가속"),
        ("2024+", "추론 특화 고도화", "단순 생성에서 복잡한 문제해결 중심의 reasoning 모델로 진화"),
    ]
    y = 1.68
    for year, title, desc in timeline:
        add_shape(slide, MSO_SHAPE.OVAL, 4.56, y + 0.14, 0.52, 0.52, ACCENT, None)
        add_card(slide, 5.18, y, 7.1, 0.92, f"{year} | {title}", desc, accent=NAVY, body_font=12, title_font=14)
        y += 1.0


def add_landscape_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "LECTURE", "2026 모델·서비스 지형도", "현재 기준 SOTA, 오픈소스, 서비스 유형, 진화 방향을 4분면으로 정리", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.65,
        5.72,
        2.0,
        "프론티어 / SOTA 모델",
        [
            "OpenAI GPT-4.5 / GPT-5(예정), o1·o3(추론 특화)",
            "Google Gemini 1.5 / 2.0 Pro",
            "Anthropic Claude 3.5 Sonnet / Opus",
            "기준 시점: 2026-04, 예시는 빠르게 바뀔 수 있음",
        ],
        accent=TEAL,
        body_font=13.5,
    )
    add_card(
        slide,
        6.8,
        1.65,
        5.83,
        2.0,
        "오픈소스 / 소형화(SLM)",
        [
            "Meta Llama 3 시리즈",
            "Mistral 등 도메인 특화 모델",
            "온디바이스·특정 기기 적용 확대",
            "소형 모델과 업무 맞춤형 배치 전략 중요",
        ],
        accent=SAGE,
        body_font=13.5,
    )
    add_card(
        slide,
        0.82,
        3.95,
        5.72,
        2.0,
        "최근 진화 흐름",
        [
            "텍스트 중심 LLM → 이미지·음성·문서를 아우르는 멀티모달",
            "질의응답 중심 → 계획하고 행동하는 Agentic AI",
            "소프트웨어 자동화 → Physical AI와의 결합",
        ],
        accent=ACCENT,
        body_font=13.5,
    )
    add_card(
        slide,
        6.8,
        3.95,
        5.83,
        2.0,
        "대표 서비스 예시",
        [
            "범용 생성형 AI: ChatGPT, Gemini, Claude",
            "AI 코딩/자동화 도구: Codex, Gemini CLI, Claude Code, GitHub Copilot",
            "핵심은 서비스 이름보다 업무 적합성과 데이터 보안 기준",
        ],
        accent=NAVY,
        body_font=13.5,
    )


def add_mechanism_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "LECTURE", "생성형 AI의 작동 원리와 HRD 접점", "모델 원리를 단순화하고 HRD 전 과정의 적용 포인트를 나란히 배치", len(prs.slides))
    add_card(slide, 0.82, 1.62, 2.85, 1.8, "LLM", "대규모 텍스트 패턴을 바탕으로 다음 토큰을 예측하는 모델", accent=TEAL, body_font=13)
    add_card(slide, 3.86, 1.62, 2.85, 1.8, "Transformer", "입력 간 관계를 병렬적으로 파악하는 구조", accent=ACCENT, body_font=13)
    add_card(slide, 0.82, 3.72, 2.85, 1.8, "Attention", "문맥 속에서 중요한 정보에 더 큰 가중치를 두는 메커니즘", accent=SAGE, body_font=13)
    add_card(slide, 3.86, 3.72, 2.85, 1.8, "멀티모달", "텍스트·이미지·음성·영상·문서를 함께 이해하고 생성", accent=GOLD, body_font=13)
    add_card(
        slide,
        7.05,
        1.62,
        5.58,
        3.9,
        "HRD와 AI의 접점",
        [
            "분석: 요구분석, 교육대상자 분류, VOC 정리",
            "설계: 과정안 작성, 메시지 구조화, 학습경험 설계",
            "개발: 교안, 슬라이드, FAQ, 웹페이지, 메일, 보고서 초안 생성",
            "운영: 문의응대, 일정공지, 후속안내, 실습자료 배포",
            "평가: 만족도 분석, 정성의견 요약, 토픽 도출, 경영진 보고자료 작성",
        ],
        accent=NAVY,
        body_font=14,
        title_font=18,
    )


def add_prompt_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "LECTURE", "프롬프트 엔지니어링 최소 구조", "출력 품질을 안정화하는 5요소와 현장 생산성 팁", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        6.0,
        4.84,
        "프롬프트 최소 구조",
        [
            "역할: 누구의 관점에서 답할지",
            "맥락: 어떤 배경과 제약이 있는지",
            "입력자료: 표, 문서, 설문, 사례 등",
            "출력형식: 표, 메일, 보고서, FAQ, HTML 등",
            "품질기준: 정확성, 톤앤매너, 길이, 금지사항",
        ],
        accent=TEAL,
        body_font=16,
        title_font=18,
    )
    add_card(
        slide,
        7.08,
        1.62,
        5.55,
        2.28,
        "생산성 팁",
        [
            "STT 활용: Windows 음성입력 `Win + H`",
            "반복 업무는 템플릿화하고, 최종 판단은 사람이 담당",
        ],
        accent=ACCENT,
        body_font=15,
        title_font=18,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.08, 4.12, 5.55, 2.34, BG_DARK, None)
    add_textbox(
        slide,
        7.36,
        4.46,
        4.95,
        1.28,
        "도구 숙련보다 중요한 것:\n업무 문제를 명확히 정의하고,\n사람이 검증 가능한 결과물로 연결하는 설계력",
        font_name=TITLE_FONT,
        font_size=18,
        color=WHITE,
        bold=True,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_practice1_strategy(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 1", "실습 1 목표와 정보 제공 방식 판단", "챗봇과 정적 웹페이지를 비교해 운영 리스크가 낮은 구성을 찾는 단계", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        4.02,
        4.85,
        "실습 목표",
        [
            "교육민원과 반복 문의를 줄이기 위한 정보 제공 방식 설계",
            "챗봇과 정적 웹페이지 중 적합한 형태 판단",
            "GitHub Pages를 활용한 교육 운영용 페이지 배포",
        ],
        accent=TEAL,
        body_font=15,
    )
    add_card(
        slide,
        5.08,
        1.62,
        3.57,
        2.24,
        "FAQ 웹페이지가 유리한 경우",
        ["정보가 자주 바뀌지 않음", "정확성이 특히 중요함", "운영 리스크를 최소화해야 함"],
        accent=NAVY,
        body_font=13.5,
    )
    add_card(
        slide,
        5.08,
        4.02,
        3.57,
        2.24,
        "GPTs / 사내 챗봇이 유리한 경우",
        ["개인화된 탐색이 중요함", "대화형 응답 경험이 필요함", "자주 변하는 질의 패턴을 흡수해야 함"],
        accent=ACCENT,
        body_font=13.5,
    )
    add_card(
        slide,
        8.9,
        1.62,
        3.72,
        4.64,
        "권장 운영 방식",
        [
            "현실적으로는 '웹페이지 + 필요 시 챗봇' 조합이 가장 안정적",
            "정확한 공지는 정적 페이지에 고정",
            "대화형 탐색은 챗봇으로 보완",
            "모바일 우선 설계와 QA가 필수",
        ],
        accent=SAGE,
        body_font=14,
    )


def add_practice1_outputs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 1", "실습 1 산출물 정의와 필수 섹션", "실습 결과물을 명확하게 정의하고, 페이지에 반드시 포함할 요소를 고정", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        4.02,
        4.9,
        "실습 결과물 정의",
        [
            "산출물: ChatGPT GPTs, 교육 안내용 마이크로사이트 1종",
            "참고 데이터: `Webpage_Sample_Data.md`",
            "배포 경로: GitHub 저장소 + GitHub Pages",
            "운영 자료에는 모바일 접속용 QR 코드 반영",
        ],
        accent=TEAL,
        body_font=14,
    )
    add_textbox(slide, 5.06, 1.68, 3.1, 0.3, "필수 섹션", font_size=17, bold=True, color=TEXT, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    pills = [
        ("과정 개요", TEAL),
        ("일정/장소", NAVY),
        ("준비물 및 사전과제", ACCENT),
        ("FAQ", SAGE),
        ("문의처", TEAL),
        ("자료 다운로드 / 링크", GOLD),
        ("모바일 접속용 QR 코드", NAVY),
    ]
    positions = [(5.06, 2.15), (7.27, 2.15), (9.48, 2.15), (5.06, 3.12), (7.27, 3.12), (9.48, 3.12), (6.17, 4.09)]
    for (label, color), (x, y) in zip(pills, positions):
        pill_w = 1.95 if x != 6.17 else 4.04
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pill_w, 0.68, color, None)
        add_textbox(slide, x + 0.14, y + 0.23, pill_w - 0.28, 0.18, label, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    add_card(
        slide,
        5.02,
        5.08,
        7.62,
        1.44,
        "운영 체크포인트",
        ["사람이 수정해야 할 항목: 오탈자, 실제 일정/장소/문의처, 사내 표현 방식, 개인정보 노출 여부"],
        accent=ACCENT,
        body_font=14,
        title_font=16,
    )


def add_practice1_workflow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 1", "실습 1 절차", "저장소 생성부터 모바일 검수와 QR 코드 반영까지 8단계로 정리", len(prs.slides))
    steps = [
        "GitHub에서 새 저장소 생성",
        "로컬에 clone",
        "웹페이지에 필요한 정보 정리",
        "AI에게 웹페이지 초안 생성 요청",
        "오탈자·일정·문의처·개인정보 검토",
        "수정본을 GitHub에 push",
        "GitHub Pages 배포 후 모바일 확인",
        "QR 코드 생성 후 현장 배포자료 삽입",
    ]
    colors = [TEAL, NAVY, SAGE, ACCENT, GOLD, TEAL, NAVY, ACCENT]
    idx = 0
    for row in range(2):
        for col in range(4):
            x = 0.82 + col * 3.12
            y = 1.75 + row * 2.2
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.72, 1.68, CARD, BORDER, 1.0)
            add_shape(slide, MSO_SHAPE.OVAL, x + 0.2, y + 0.2, 0.56, 0.56, colors[idx], None)
            add_textbox(slide, x + 0.39, y + 0.38, 0.18, 0.12, str(idx + 1), font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
            add_textbox(slide, x + 0.92, y + 0.18, 1.55, 1.1, steps[idx], font_size=14, color=TEXT, bold=True, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
            idx += 1


def add_practice1_prompt(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 1", "실습 프롬프트 예시", "웹사이트 생성 요청을 구조화된 요구조건 형태로 제시", len(prs.slides))
    prompt = """HRD 과정 안내용 단일 페이지 웹사이트를 만들어줘.

목표:
- 교육 참가자가 일정, 장소, 준비물, FAQ, 문의처를 빠르게 확인할 수 있게 하기
- 모바일에서 먼저 읽히도록 설계하기
- 대상자명단에서는 전화걸기, 문자보내기 버튼이 있어야함

반드시 포함할 섹션:
- 과정 개요
- 일정/장소
- 준비물
- FAQ 6개
- 문의처
- 자료 다운로드 버튼

디자인 조건:
- 제조업 대기업의 신뢰감 있는 톤
- 과도한 애니메이션 금지
- 가독성 우선
- HTML/CSS/JavaScript를 분리해서 작성"""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 1.68, 8.15, 5.2, BG_DARK, None)
    add_textbox(slide, 1.04, 1.98, 7.7, 4.72, prompt, font_name=MONO_FONT, font_size=11.4, color=WHITE, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    add_card(
        slide,
        9.2,
        1.68,
        3.42,
        2.38,
        "프롬프트 포인트",
        [
            "목표, 필수 섹션, 디자인 조건을 분리",
            "대상자명단의 전화/문자 CTA를 명시",
            "모바일 우선과 파일 분리까지 요구",
        ],
        accent=ACCENT,
        body_font=13.5,
    )
    add_card(
        slide,
        9.2,
        4.32,
        3.42,
        2.28,
        "현업 검토 포인트",
        [
            "정확한 일정·장소·문의처 반영 여부",
            "사내 문체와 보안·개인정보 기준 충족 여부",
            "모바일 화면에서 실제로 읽히는지 확인",
        ],
        accent=TEAL,
        body_font=13.2,
    )


def add_demo_dataset_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "DEMO 1", "시연 1 목표와 데이터 구조", "리더십 교육 F/U 개인화 메시지 자동화를 위한 입력 스키마 정의", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        3.8,
        4.95,
        "시연 목표",
        [
            "Action Plan 데이터를 활용해 개인화된 후속 메시지 초안을 자동 생성",
            "메일 작성 시간을 줄이면서도 개인별 맥락을 유지",
            "개인화된 메일을 발송하는 Tool 활용 흐름 제시",
        ],
        accent=TEAL,
        body_font=15,
    )
    schema_cards = [
        ("기본 정보", "`name`, `company`, `department`, `position`", NAVY),
        ("과정 산출물", "`aspiration`, `action_plan`, `expected_behavior`", ACCENT),
        ("피드백 정보", "`training_comment`, `top_takeaway`", SAGE),
        ("운영 정보", "`email`, `manager_name`", GOLD),
    ]
    positions = [(4.9, 1.62), (9.0, 1.62), (4.9, 4.02), (9.0, 4.02)]
    for (title, body, color), (x, y) in zip(schema_cards, positions):
        add_card(slide, x, y, 3.58, 2.05, title, body, accent=color, body_font=12.4, title_font=16)


def add_demo_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "DEMO 1", "권장 데이터 흐름", "CSV 로드부터 HTML 메일 생성과 일괄 발송까지 6단계 자동화", len(prs.slides))
    steps = [
        ("1. 로드/전처리", "스프레드시트 데이터(`교육_액션플랜_데이터.csv`) 로드 및 전처리"),
        ("2. 프롬프트 템플릿", "이름, 소속, 실행계획, 코멘트, 매니저명 등을 변수로 매핑"),
        ("3. 메시지 생성", "지정된 담당 매니저 명의의 개인별 후속지원 메시지 초안 생성"),
        ("4. HTML 변환", "HTML 메일 본문 형식으로 변환"),
        ("5. 최종 검수", "어조, 민감표현, 인사말 등 사실관계 검수"),
        ("6. 일괄 발송", "XLMultimail 도구를 활용해 전체 발송"),
    ]
    colors = [TEAL, NAVY, ACCENT, SAGE, GOLD, NAVY]
    for idx, ((title, body), color) in enumerate(zip(steps, colors)):
        row = idx // 3
        col = idx % 3
        x = 0.82 + col * 4.12
        y = 1.82 + row * 2.18
        add_card(slide, x, y, 3.68, 1.78, title, body, accent=color, body_font=12.6, title_font=15)
    for x1, x2, y in [(4.5, 4.94, 2.72), (8.62, 9.06, 2.72), (4.5, 4.94, 4.9), (8.62, 9.06, 4.9)]:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, i(x1), i(y), i(x2), i(y))
        connector.line.color.rgb = rgb(BORDER)
        connector.line.width = pt(1.4)


def add_demo_tools_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "DEMO 1", "시연 도구와 파일", "자동화 스크립트, 데이터셋, 발송 도구의 역할을 분리해 보여주기", len(prs.slides))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 1.76, 11.8, 3.52, CARD, BORDER, 1.0)
    table = slide.shapes.add_table(4, 3, i(1.0), i(2.0), i(11.4), i(3.04)).table
    headers = ["구분", "파일명", "설명"]
    widths = [1.5, 4.6, 5.3]
    for idx, width in enumerate(widths):
        table.columns[idx].width = i(width)
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BG_DARK)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(WHITE)
    rows = [
        ("자동화 스크립트", "SCRIPT_glm_feedback_automation.py", "CSV 데이터를 읽어 LLM API로 개인별 피드백 초안을 생성하고 HTML 메일 미리보기와 대시보드까지 출력"),
        ("실습 데이터", "sample_data/교육_액션플랜_데이터.csv", "50명 규모 가상 수강자의 액션플랜·피드백 데이터"),
        ("메일 발송 도구", "FILE_엑셀에서메일발송_XLMultimail_v2.6.0_20240117.zip", "Excel 기반 일괄 메일 발송 애드인(XLMultimail)"),
    ]
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(CARD)
            cell.text = wrap_file(value, 24 if c_idx == 1 else 48)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = MONO_FONT if c_idx == 1 else BODY_FONT
            p.runs[0].font.size = pt(11.5)
            p.runs[0].font.color.rgb = rgb(TEXT)
    add_card(
        slide,
        0.82,
        5.58,
        11.8,
        0.92,
        "발송 전 최종 검수",
        ["어조, 민감표현, 인사말, 사실관계는 사람이 마지막으로 확인해야 함"],
        accent=ACCENT,
        body_font=13.5,
        title_font=16,
    )


def add_practice2_data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 2", "실습 2 목표와 데이터 프레이밍", "정량·정성 분석과 HTML 보고서 산출을 한 흐름으로 묶기", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        4.0,
        4.96,
        "실습 목표",
        [
            "교육 만족도 데이터를 정량/정성으로 나누어 분석 흐름 이해",
            "수작업 보고를 줄이고 반복 가능한 분석 파이프라인 설계",
            "결과를 담은 데이터 대시보드를 단일 HTML 보고서로 생성",
        ],
        accent=TEAL,
        body_font=15,
    )
    add_card(
        slide,
        5.1,
        1.62,
        3.56,
        2.26,
        "정량 데이터",
        ["리커트 척도 기반 평균(Average)", "추천의향 문항은 NPS 별도 계산", "강사만족도, 내용적합성 등 문항군 분석"],
        accent=NAVY,
        body_font=13.2,
    )
    add_card(
        slide,
        8.96,
        1.62,
        3.68,
        2.26,
        "정성 데이터",
        ["주관식 수강소감에 대한 KOTE 감성분석", "핵심 토픽모델링으로 대표 인사이트 도출", "대표 키워드·감성분포 시각화"],
        accent=ACCENT,
        body_font=13.2,
    )
    add_card(
        slide,
        5.1,
        4.16,
        7.54,
        2.42,
        "2차 산출물",
        [
            "NPS 차트, 감성분포, 토픽모델링 워드클라우드가 결합된 HTML 기반 최종 요약 보고서",
            "경영진 보고용으로 정량 수치와 정성 인사이트를 한 화면에 통합",
            "후속 액션 도출까지 연결되는 환류(Feedback Loop) 설계",
        ],
        accent=SAGE,
        body_font=14,
        title_font=18,
    )


def add_practice2_pipeline_slide(prs: Presentation, title: str, rows, footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 2", title, "수집 → 전처리 → 정량/정성 분석 → HTML 생성 → 환류의 파이프라인", len(prs.slides))
    headers = ["단계", "입력", "처리", "출력"]
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 1.72, 11.8, 4.72, CARD, BORDER, 1.0)
    table = slide.shapes.add_table(len(rows) + 1, 4, i(0.98), i(1.92), i(11.45), i(4.3)).table
    widths = [1.25, 1.8, 4.95, 3.45]
    for idx, width in enumerate(widths):
        table.columns[idx].width = i(width)
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BG_DARK)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = pt(11.5)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(WHITE)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(CARD)
            cell.text = wrap_file(value, 22 if c_idx == 0 else 24 if c_idx == 1 else 40)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = BODY_FONT
            p.runs[0].font.size = pt(10.2)
            p.runs[0].font.color.rgb = rgb(TEXT)
    if footer:
        add_textbox(slide, 0.92, 6.55, 11.2, 0.28, footer, font_size=11, color=MUTED, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)


def add_practice2_tools_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "PRACTICE 2", "실습 2 결과물, 도구, 참고자료", "최종 HTML 보고서와 이를 지원하는 데이터·프롬프트·논문 묶음을 정리", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        4.05,
        4.92,
        "실무형 결과물 예시",
        ["KOTE 감성분석 요약표", "문항별 NPS 지표", "토픽모델링 워드클라우드", "정량·정성 결과를 하나로 렌더링한 HTML 결과 보고서"],
        accent=TEAL,
        body_font=15,
    )
    add_card(
        slide,
        5.12,
        1.62,
        7.52,
        3.08,
        "실습 파일",
        [
            "`sample_data/교육만족도_설문_데이터.csv` : 300명 규모 만족도 데이터",
            "`PROMPT_Quantitative_Analysis_Prompt_Sample.md` : 평균·NPS 요청 프롬프트 예시",
            "`PROMPT_TopicModeling_Prompt_Sample.md` : 감성분석·토픽모델링 요청 프롬프트 예시",
            "`REFERENCE_Mu. Y., (2024)...pdf` : LLM 기반 토픽모델링 방법론 참고 논문",
        ],
        accent=ACCENT,
        body_font=13.4,
    )
    add_card(
        slide,
        5.12,
        4.96,
        7.52,
        1.58,
        "참고 링크",
        ["KOTE 레포지토리: https://github.com/searle-j/KOTE", "Mu(2024) 논문: https://arxiv.org/pdf/2403.16248"],
        accent=NAVY,
        body_font=13.4,
        title_font=16,
    )


def add_case_story_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "CASE", "리더십 제고 사례 스토리라인", "2025년의 공감대 형성에서 2026년의 실제 문제해결·확산으로 연결", len(prs.slides))
    add_card(
        slide,
        0.92,
        1.86,
        5.28,
        4.72,
        "2025년: 중요성 인식과 공감대 형성의 해",
        ["임원: 1박 2일 LLM 실습과 산하조직 DX 계획 수립", "팀장: CTO 특강 → LLM 실습 → 과제도출 워크숍"],
        accent=NAVY,
        body_font=17,
        title_font=18,
    )
    add_shape(slide, MSO_SHAPE.CHEVRON, 6.36, 3.05, 0.9, 1.2, GOLD, None)
    add_card(
        slide,
        7.44,
        1.86,
        5.28,
        4.72,
        "2026년: 실제 문제해결과 확산의 해",
        ["Change Agent 선발", "과제 확정", "장기 수행", "성과공유 및 차년도 승계"],
        accent=ACCENT,
        body_font=17,
        title_font=18,
    )


def add_change_agent_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "CASE", "Change Agent 설명 보강", "현장 실행 담당자의 정의와 기대 역할을 명확히 제시", len(prs.slides))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 1.68, 11.82, 1.04, BG_DARK, None)
    add_textbox(slide, 1.08, 1.98, 11.2, 0.3, "정의: 실 단위 조직의 문화/업무 개선을 주도하는 현장 실행 담당자", font_size=18, color=WHITE, bold=True, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    roles = [
        ("현업 문제 발굴", "실제 현장에서 반복되거나 병목이 되는 과제를 정의"),
        ("AI 적용 실험", "작은 자동화와 파일럿을 빠르게 실행"),
        ("구성원 참여 유도", "현업 구성원에게 참여 맥락과 실익을 설명"),
        ("성과 사례 문서화", "배운 점과 결과를 재사용 가능한 사례로 축적"),
    ]
    x = 0.82
    colors = [TEAL, NAVY, ACCENT, SAGE]
    for (title, body), color in zip(roles, colors):
        add_card(slide, x, 3.18, 2.84, 2.72, title, body, accent=color, body_font=13.3, title_font=16)
        x += 3.0


def add_wrap_recommendation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "WRAP-UP", "실천 제언과 추천 탐색 영역", "학습 습관, 팀 자산화, 보안 검토, 도구 탐색을 동시에 제안", len(prs.slides))
    add_card(
        slide,
        0.82,
        1.62,
        5.35,
        4.94,
        "실천 제언",
        [
            "트렌드를 모두 따라갈 필요는 없지만, 업무와 연결되는 변화는 놓치지 않는다.",
            "매주 30분은 새로운 도구를 탐색하고, 매월 1개는 실제 업무에 적용해본다.",
            "팀 단위로 '잘 쓴 프롬프트', '잘 안 된 프롬프트', '재사용 가능한 템플릿'을 축적한다.",
            "외부 서비스 사용 시 보안과 저작권 검토를 습관화한다.",
        ],
        accent=TEAL,
        body_font=14.6,
    )
    add_card(
        slide,
        6.42,
        1.62,
        6.2,
        3.08,
        "추천 탐색 영역",
        [
            "범용 생성형 AI: ChatGPT, Gemini, Claude",
            "에이전트형 업무도구: Codex, Gemini CLI, Claude Code, GitHub Copilot",
            "검색/리서치형 AI 서비스: 사내 보안정책과 업무 적합성 기준으로 선별 도입",
        ],
        accent=ACCENT,
        body_font=14.2,
    )
    add_card(
        slide,
        6.42,
        4.98,
        6.2,
        1.58,
        "학습 채널",
        ["GeekNews", "유튜브 기술 채널", "실무 중심 도서"],
        accent=NAVY,
        body_font=14.2,
        title_font=16,
    )


def add_utility_slide(prs: Presentation, title: str, intro: str, features, accent_color: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "UTILITY", title, intro, len(prs.slides))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 1.62, 11.82, 0.86, BG_DARK, None)
    add_textbox(slide, 1.08, 1.92, 11.2, 0.26, intro, font_size=14, color=WHITE, margin_left=0, margin_right=0, margin_top=0, margin_bottom=0)
    positions = [(0.82, 2.82), (6.02, 2.82), (0.82, 5.02), (6.02, 5.02)]
    for (name, body), (x, y) in zip(features, positions):
        add_card(slide, x, y, 5.25, 1.76, name, body, accent=accent_color, body_font=13.2, title_font=16)


def add_filemap_slide(prs: Presentation, title: str, rows) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_LIGHT)
    add_content_title(slide, "APPENDIX", title, "프로젝트 폴더(`LETSAI/`)의 파일과 관련 세션을 연결한 부록", len(prs.slides))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 1.62, 11.95, 4.95, CARD, BORDER, 1.0)
    table = slide.shapes.add_table(len(rows) + 1, 4, i(0.88), i(1.88), i(11.6), i(4.45)).table
    headers = ["파일/폴더", "유형", "관련 세션", "설명"]
    widths = [3.15, 1.1, 1.8, 5.55]
    for idx, width in enumerate(widths):
        table.columns[idx].width = i(width)
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BG_DARK)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = pt(11.2)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(WHITE)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(CARD)
            cell.text = wrap_file(value, 28 if c_idx == 0 else 14 if c_idx == 1 else 18 if c_idx == 2 else 42)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = MONO_FONT if c_idx == 0 else BODY_FONT
            p.runs[0].font.size = pt(9.4)
            p.runs[0].font.color.rgb = rgb(TEXT)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = i(SLIDE_W)
    prs.slide_height = i(SLIDE_H)

    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "HRD부문의 효과적인 AI 활용 방안"
    prs.core_properties.subject = "현대제철 HRD 세미나용 구조화 발표자료"
    prs.core_properties.language = "ko-KR"
    prs.core_properties.keywords = "HRD, AI, 현대제철, 교육운영, 발표자료"

    add_cover_slide(prs)
    add_overview_slide(prs)
    add_core_message_slide(prs)
    add_agenda_slide(prs)

    add_divider_slide(prs, "01", "강의", "생성형 AI 동향 이해와 HRD 적용 관점 정리")
    add_lecture_goal_timeline(prs)
    add_landscape_slide(prs)
    add_mechanism_slide(prs)
    add_prompt_slide(prs)

    add_divider_slide(prs, "02", "실습 1", "교육운영 효율화를 위한 AI 챗봇 / 웹페이지 만들기")
    add_practice1_strategy(prs)
    add_practice1_outputs(prs)
    add_practice1_workflow(prs)
    add_practice1_prompt(prs)

    add_divider_slide(prs, "03", "시연 1", "교육 콘텐츠 및 피드백 자동화")
    add_demo_dataset_slide(prs)
    add_demo_flow_slide(prs)
    add_demo_tools_slide(prs)

    add_divider_slide(prs, "04", "실습 2", "교육 데이터의 AI 기반 자동화 전략")
    add_practice2_data_slide(prs)
    practice2_rows = [
        ("수집", "설문 원본", "결측치 및 항목 속성 점검 (`교육만족도_설문_데이터.csv`)", "분석 대상 데이터셋"),
        ("전처리", "정량/정성 응답", "KOTE 감성분석을 위한 형태소/텍스트 정제, 비식별화", "정제된 분석용 데이터"),
        ("정량분석", "점수형 응답", "만족도 문항군 평균 도출, 추천도 기반 NPS 계산", "주요 통계 요약 및 차트 값"),
        ("정성분석", "서술형 응답('수강소감')", "KOTE 감성분류, 토픽모델링(LDA 등)으로 주요 키워드 추출", "텍스트 분석 핵심 인사이트"),
        ("HTML 생성", "정량/정성 결과 종합", "시각화를 포함해 하나의 문서로 결합", "경영진 보고용 HTML 통합 보고서"),
        ("환류", "최종 요약 보고서", "데이터에 나타난 강약점을 기반으로 다음 차수 액션 도출", "교육 개선 실행 과제"),
    ]
    add_practice2_pipeline_slide(prs, "권장 분석 흐름 (1/2)", practice2_rows[:3])
    add_practice2_pipeline_slide(
        prs,
        "권장 분석 흐름 (2/2)",
        practice2_rows[3:],
        "참고: KOTE 레포지토리와 Mu(2024) 논문을 통해 모델 다운로드 및 LLM 기반 토픽모델링 가이드를 확인",
    )
    add_practice2_tools_slide(prs)

    add_divider_slide(prs, "05", "사례공유", "AI 기반 혁신적 성과창출을 위한 리더십 제고 사례")
    add_case_story_slide(prs)
    add_change_agent_slide(prs)

    add_divider_slide(prs, "06", "Wrap-up", "HRD 담당자의 AI 학습방법 및 전사 AX 추진 제언")
    add_wrap_recommendation_slide(prs)
    add_utility_slide(
        prs,
        "추천 Utility: PowerToys",
        "Windows 생산성을 높여주는 무료 유틸리티 모음으로, 실습 중 창 정리와 빠른 실행, 텍스트 추출에 특히 유용",
        [
            ("FancyZones", "여러 창을 미리 정한 레이아웃에 맞춰 배치해 강의자료, AI 도구, 메모 창을 동시에 안정적으로 운영"),
            ("Text Extractor", "이미지·영상·웹페이지에 보이는 텍스트를 OCR로 바로 복사해 예시 화면 문구를 재활용"),
            ("PowerToys Run / Command Palette", "앱, 파일, 폴더, 명령을 빠르게 실행해 실습 중 도구 전환 속도를 높임"),
            ("Always On Top", "체크리스트나 공지사항 창을 항상 위에 고정해 진행 안정성을 높임"),
        ],
        TEAL,
    )
    add_utility_slide(
        prs,
        "추천 Utility: ShareX",
        "화면 캡처, 스크롤 캡처, GIF 녹화, OCR, 주석 편집까지 한 번에 처리할 수 있어 교육 운영 자료 제작에 매우 실용적",
        [
            ("Region Capture / Scrolling Capture", "긴 웹페이지나 설문 결과 화면 전체를 한 장으로 캡처해 FAQ 페이지와 대시보드 공유에 활용"),
            ("Annotation / Image Editor", "화살표, 박스, 흐림, 하이라이트 등을 넣어 매뉴얼·안내 이미지를 빠르게 제작"),
            ("Screen Recording / GIF", "짧은 사용법 시연을 GIF 또는 영상으로 남겨 교육 전후 안내자료로 전환"),
            ("OCR / Workflow", "캡처한 이미지에서 텍스트를 추출하고, 저장·복사·업로드 등 후속 작업을 자동화"),
        ],
        ACCENT,
    )

    filemap_rows = [
        ("DOCUMENT_Design.md", "문서", "전체", "강의 및 실습 설계 문서(본 문서)"),
        ("DOCUMENT_4월 HRD부문의 AI 활용 세미나 (1차명단).pdf", "문서", "전체", "참석자 1차 확정 명단"),
        ("Webpage_Sample_Data.md", "데이터", "실습 1 (§5)", "교육 안내 웹페이지 제작용 가상 샘플 정보"),
        ("SCRIPT_glm_feedback_automation.py", "스크립트", "시연 1 (§6)", "CSV → LLM API → 개인화 피드백 HTML 메일 자동 생성 파이프라인"),
        ("FILE_엑셀에서메일발송_XLMultimail_v2.6.0_20240117.zip", "도구", "시연 1 (§6)", "Excel 기반 일괄 메일 발송 애드인(XLMultimail)"),
        ("PROMPT_Quantitative_Analysis_Prompt_Sample.md", "프롬프트", "실습 2 (§7)", "정량 분석(평균·NPS) 요청용 프롬프트 예시"),
        ("PROMPT_TopicModeling_Prompt_Sample.md", "프롬프트", "실습 2 (§7)", "감성분석·토픽모델링 요청용 프롬프트 예시"),
        ("REFERENCE_Mu. Y., (2024) Large Language Models Offer an Alternative to the Traditional.pdf", "참고문헌", "실습 2 (§7)", "LLM 기반 토픽모델링 방법론 논문"),
        ("sample_data/교육_액션플랜_데이터.csv", "데이터", "시연 1 (§6)", "50명 가상 수강자의 액션플랜·피드백 원본 데이터"),
        ("sample_data/교육만족도_설문_데이터.csv", "데이터", "실습 2 (§7)", "300명 교육 만족도 설문 데이터(정량 + 정성)"),
        ("sample_data/course_info.csv", "데이터", "실습 1 (§5)", "과정 기본 정보"),
        ("sample_data/course_faq.csv", "데이터", "실습 1 (§5)", "과정 FAQ 데이터"),
        ("config/", "설정", "전체", "스크립트 실행 환경 설정"),
    ]
    add_filemap_slide(prs, "부록: 폴더 파일맵 (1/3)", filemap_rows[:5])
    add_filemap_slide(prs, "부록: 폴더 파일맵 (2/3)", filemap_rows[5:9])
    add_filemap_slide(prs, "부록: 폴더 파일맵 (3/3)", filemap_rows[9:])

    OUTPUT_DIR.mkdir(exist_ok=True)
    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    output_path = build_presentation()
    print(output_path)
